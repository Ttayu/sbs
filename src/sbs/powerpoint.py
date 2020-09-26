import io
from datetime import datetime
from enum import IntEnum, auto
from pathlib import Path
from typing import Union

import pptx
from PIL import Image
from pptx.presentation import Presentation

TEMPLATE = Path(__file__).parents[2] / "source" / "./template.pptx"


class SlideMaster(IntEnum):
    Title = 0
    Image = auto()
    Contents = auto()
    ThreeComparisons = auto()
    TwoContents = auto()
    TwoComparisons = auto()
    TitleOnly = auto()
    Blank = auto()


ImageFileType = Union[str, io.BytesIO, Path]


class PowerPoint(Presentation):
    def __init__(self):
        prs = pptx.Presentation(TEMPLATE)
        super().__init__(prs.element, prs.part)
        self.last_slide = SlideMaster.Title

    def save(self, file: Union[str, Path]) -> None:
        if Path(file).suffix != ".pptx":
            file = str(file) + ".pptx"
        super().save(str(file))

    def add_slide(self, slide: SlideMaster):
        slide_layout: pptx.slide.SlideLayout = self.slide_layouts[slide]
        self.slides.add_slide(slide_layout)
        self.last_slide = slide

    def add_title(self, title: str = "title", name: str = "") -> None:
        self.add_slide(SlideMaster.Title)
        self.update_title(title, name)

    def update_title(self, title: str = "title", name: str = "") -> None:
        assert self.last_slide == SlideMaster.Title
        self.slides[-1].placeholders[0].text = title
        self.slides[-1].placeholders[1].text = (
            name if name else str(datetime.now().date())
        )

    def add_contents(self, title: str = "title", contents: str = "") -> None:
        self.add_slide(SlideMaster.Contents)
        self.update_title(title, contents)

    def update_contents(self, title: str = "title", contents: str = "") -> None:
        assert self.last_slide == SlideMaster.Contents
        self.slides[-1].placeholders[0].text = title
        self.slides[-1].placeholders[1].text = contents

    def add_image(self, image: Union[Image.Image, ImageFileType]) -> None:
        self.add_slide(SlideMaster.Image)
        self.update_image(image)

    def update_image(self, image: Union[Image.Image, ImageFileType]) -> None:
        assert self.last_slide == SlideMaster.Image
        self._plot_image(image)

    def _save_inmemory_PIL(self, image: Image.Image) -> io.BytesIO:
        item = io.BytesIO()
        image.save(item, "png")
        item.seek(0)
        return item

    def _plot_image(self, image: Union[Image.Image, Path]) -> None:
        if isinstance(image, Image.Image):
            item: Union[str, io.BytesIO] = self._save_inmemory_PIL(image)
        else:
            item = str(image)

        pixture_placeholder: pptx.shapes.placeholder.PicturePlaceholder = self.slides[
            -1
        ].placeholders[13]
        assert isinstance(
            pixture_placeholder, pptx.shapes.placeholder.PicturePlaceholder
        )

        # Insert the picture
        placeholder: pptx.shapes.placeholder.PlaceholderPicture = (
            pixture_placeholder.insert_picture(item)
        )

        # Calculate the image size of the image
        width, height = placeholder.image.size

        # Calculate ratios and compare
        image_ratio = width / height
        placeholder_ratio = placeholder.width / placeholder.height
        ratio_difference = placeholder_ratio - image_ratio

        placeholder.crop_top = 0
        placeholder.crop_left = 0
        placeholder.crop_bottom = 0
        placeholder.crop_right = 0

        if ratio_difference > 0:
            width_slide = placeholder.width
            placeholder.height = placeholder.height
            placeholder.width = int((placeholder.height) * image_ratio)
            placeholder.left = (width_slide - placeholder.width) // 2
        else:
            height_slide = placeholder.height
            placeholder.width = placeholder.width
            placeholder.height = int(placeholder.width / image_ratio)
            placeholder.top = (height_slide - placeholder.height) // 2
        assert placeholder.width > 0
        assert placeholder.height > 0
        assert placeholder.top > 0 or placeholder.left > 0

    def _current_placeholder_id(self):
        for shape in self.slides[-1].shapes:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                print(phf.idx, phf.type)


def paste_image(
    input_dir: Union[str, Path],
    output_file: Union[str, Path],
    title="",
    glob_pattern: str = "*.png",
) -> None:
    ppt = PowerPoint()
    ppt.update_title(title)
    for file in Path(input_dir).glob(glob_pattern):
        ppt.add_image(file)
    ppt.save(output_file)
